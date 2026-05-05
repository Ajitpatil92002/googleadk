import asyncio
import json
import os
from typing import Optional

from fpdf import FPDF
from google.adk.tools.function_tool import FunctionTool
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


async def create_pdf(
    title: str,
    sections_json: str,
    output_path: str,
    author: Optional[str] = None,
) -> dict:
    """
    Create a PDF document with text, tables, and multiple sections.

    Args:
        title: The document title shown on the first page.
        sections_json: A JSON string representing a list of sections. Each section is an object with:
            - "heading" (string): Section heading text.
            - "paragraphs" (list of strings): Body text paragraphs.
            - "table" (object): Optional table with "headers" (list of strings) and "rows" (list of lists of strings).
            Example: '[{"heading": "Summary", "paragraphs": ["This is a test.", "Another paragraph."], "table": {"headers": ["Name", "Value"], "rows": [["A", "1"], ["B", "2"]]}}]'
        output_path: File path to save the PDF (e.g., 'report.pdf').
        author: Optional author name for the PDF metadata.

    Returns:
        Dictionary with 'status', 'file_path', and 'message'.
    """
    try:
        sections = json.loads(sections_json)
    except json.JSONDecodeError as e:
        return {"status": "error", "message": f"Invalid JSON in sections_json: {e}"}

    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        if author:
            pdf.set_author(author)
        pdf.set_title(title)

        pdf.add_page()
        pdf.set_font("Helvetica", "B", 24)
        pdf.cell(0, 20, title, ln=True, align="C")
        pdf.ln(10)

        for section in sections:
            heading = section.get("heading", "")
            paragraphs = section.get("paragraphs", [])
            table = section.get("table")

            if heading:
                pdf.set_font("Helvetica", "B", 16)
                pdf.cell(0, 10, heading, ln=True)
                pdf.ln(2)

            if paragraphs:
                pdf.set_font("Helvetica", "", 11)
                for p in paragraphs:
                    pdf.multi_cell(0, 6, str(p))
                    pdf.ln(3)

            if table:
                headers = table.get("headers", [])
                rows = table.get("rows", [])

                if headers:
                    col_count = len(headers)
                    available_width = pdf.w - 20
                    col_width = available_width / col_count

                    pdf.set_font("Helvetica", "B", 10)
                    pdf.set_fill_color(66, 133, 244)
                    pdf.set_text_color(255, 255, 255)
                    for h in headers:
                        pdf.cell(col_width, 8, str(h), border=1, fill=True, align="C")
                    pdf.ln()

                    pdf.set_font("Helvetica", "", 10)
                    pdf.set_text_color(0, 0, 0)
                    for i, row in enumerate(rows):
                        if i % 2 == 0:
                            pdf.set_fill_color(240, 240, 240)
                        else:
                            pdf.set_fill_color(255, 255, 255)
                        for cell in row:
                            pdf.cell(col_width, 7, str(cell), border=1, fill=True)
                    pdf.ln()

                pdf.ln(5)

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        pdf.output(output_path)
        abs_path = os.path.abspath(output_path)
        return {
            "status": "success",
            "file_path": abs_path,
            "message": f"PDF created successfully at {abs_path}",
        }
    except Exception as e:
        return {"status": "error", "message": f"Failed to create PDF: {str(e)}"}


_XL_THEMES = [
    {"name": "Blue",     "primary": "1A73E8", "header_bg": "1A73E8", "header_font": "FFFFFF", "alt_row": "E8F0FE", "border": "BDC3C7", "title_bg": "0D47A1", "title_font": "FFFFFF", "sub_bg": "42A5F5", "sub_font": "FFFFFF"},
    {"name": "Green",    "primary": "2E7D32", "header_bg": "2E7D32", "header_font": "FFFFFF", "alt_row": "E8F5E9", "border": "A5D6A7", "title_bg": "1B5E20", "title_font": "FFFFFF", "sub_bg": "66BB6A", "sub_font": "FFFFFF"},
    {"name": "Purple",   "primary": "6A1B9A", "header_bg": "6A1B9A", "header_font": "FFFFFF", "alt_row": "F3E5F5", "border": "CE93D8", "title_bg": "4A148C", "title_font": "FFFFFF", "sub_bg": "AB47BC", "sub_font": "FFFFFF"},
    {"name": "Orange",   "primary": "E65100", "header_bg": "E65100", "header_font": "FFFFFF", "alt_row": "FFF3E0", "border": "FFB74D", "title_bg": "BF360C", "title_font": "FFFFFF", "sub_bg": "FF7043", "sub_font": "FFFFFF"},
    {"name": "Teal",     "primary": "00695C", "header_bg": "00695C", "header_font": "FFFFFF", "alt_row": "E0F2F1", "border": "80CBC4", "title_bg": "004D40", "title_font": "FFFFFF", "sub_bg": "26A69A", "sub_font": "FFFFFF"},
    {"name": "Red",      "primary": "C62828", "header_bg": "C62828", "header_font": "FFFFFF", "alt_row": "FFEBEE", "border": "EF9A9A", "title_bg": "B71C1C", "title_font": "FFFFFF", "sub_bg": "EF5350", "sub_font": "FFFFFF"},
    {"name": "Dark",     "primary": "37474F", "header_bg": "37474F", "header_font": "FFFFFF", "alt_row": "ECEFF1", "border": "90A4AE", "title_bg": "263238", "title_font": "FFFFFF", "sub_bg": "607D8B", "sub_font": "FFFFFF"},
    {"name": "Rose",     "primary": "AD1457", "header_bg": "AD1457", "header_font": "FFFFFF", "alt_row": "FCE4EC", "border": "F48FB1", "title_bg": "880E4F", "title_font": "FFFFFF", "sub_bg": "EC407A", "sub_font": "FFFFFF"},
]


def _xl_apply_border(cell, border_color):
    thin = Side(style="thin", color=border_color)
    cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)


def _xl_style_range(ws, row, col_start, col_end, font=None, fill=None, alignment=None, border_color=None):
    for c in range(col_start, col_end + 1):
        cell = ws.cell(row=row, column=c)
        if font:
            cell.font = font
        if fill:
            cell.fill = fill
        if alignment:
            cell.alignment = alignment
        if border_color:
            _xl_apply_border(cell, border_color)


async def create_excel(
    sheets_json: str,
    output_path: str,
    theme: Optional[str] = None,
) -> dict:
    """
    Create a colorful, professionally styled Excel spreadsheet with support for headers, merged titles, subtitles, and themed formatting.

    Args:
        sheets_json: A JSON string representing a list of sheets. Each sheet is an object with:
            - "name" (string): Sheet tab name.
            - "title" (string): Optional large merged title row (e.g., college name). Spans all columns.
            - "subtitle" (string): Optional merged subtitle row (e.g., class name, department). Spans all columns.
            - "headers" (list of strings): Column headers.
            - "rows" (list of lists): Row data. Each cell can be a string, number, or formula string starting with '='.
            - "column_widths" (list of integers): Optional column widths.
            - "auto_fit" (bool): Optional, auto-adjust column widths. Defaults to true.
            - "merged_cells" (list of objects): Optional list of merges, each with "start_row", "start_col", "end_row", "end_col" (1-indexed).
            - "highlight_rows" (list of integers): Optional 1-indexed row numbers to highlight (relative to data rows after headers).
            Example: '{"name": "Timetable", "title": "ABC College", "subtitle": "CS - Semester 5 Timetable", "headers": ["Time", "Monday", "Tuesday", "Wednesday", "Thursday", "Friday"], "rows": [["9:00-10:00", "Math", "Physics", "Chemistry", "Math", "Lab"]]}'
        output_path: File path to save the Excel file (e.g., 'data.xlsx').
        theme: Optional theme name. Choices: Blue, Green, Purple, Orange, Teal, Red, Dark, Rose. Auto-selected if not provided.

    Returns:
        Dictionary with 'status', 'file_path', 'theme', and 'message'.
    """
    try:
        sheets = json.loads(sheets_json)
    except json.JSONDecodeError as e:
        return {"status": "error", "message": f"Invalid JSON in sheets_json: {e}"}

    try:
        import random
        chosen = None
        if theme:
            for t in _XL_THEMES:
                if t["name"].lower() == theme.lower():
                    chosen = t
                    break
        if not chosen:
            chosen = random.choice(_XL_THEMES)

        wb = Workbook()
        first = True

        for sheet_data in sheets:
            name = sheet_data.get("name", "Sheet1")
            title = sheet_data.get("title", "")
            subtitle = sheet_data.get("subtitle", "")
            headers = sheet_data.get("headers", [])
            rows = sheet_data.get("rows", [])
            column_widths = sheet_data.get("column_widths", [])
            auto_fit = sheet_data.get("auto_fit", True)
            merged_cells = sheet_data.get("merged_cells", [])
            highlight_rows = sheet_data.get("highlight_rows", [])

            if first:
                ws = wb.active
                ws.title = name
                first = False
            else:
                ws = wb.create_sheet(title=name)

            col_count = max(len(headers), 1)
            current_row = 1

            title_font = Font(bold=True, color=chosen["title_font"], size=16, name="Calibri")
            title_fill = PatternFill(start_color=chosen["title_bg"], end_color=chosen["title_bg"], fill_type="solid")
            title_align = Alignment(horizontal="center", vertical="center")

            sub_font = Font(bold=True, color=chosen["sub_font"], size=12, name="Calibri")
            sub_fill = PatternFill(start_color=chosen["sub_bg"], end_color=chosen["sub_bg"], fill_type="solid")

            header_font = Font(bold=True, color=chosen["header_font"], size=11, name="Calibri")
            header_fill = PatternFill(start_color=chosen["header_bg"], end_color=chosen["header_bg"], fill_type="solid")
            header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)

            body_font = Font(size=11, name="Calibri", color="333333")
            body_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
            alt_fill = PatternFill(start_color=chosen["alt_row"], end_color=chosen["alt_row"], fill_type="solid")
            white_fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")
            highlight_fill = PatternFill(start_color=chosen["primary"], end_color=chosen["primary"], fill_type="solid")
            highlight_font = Font(bold=True, color="FFFFFF", size=11, name="Calibri")

            if title:
                ws.merge_cells(start_row=current_row, start_column=1, end_row=current_row, end_column=col_count)
                cell = ws.cell(row=current_row, column=1, value=title)
                cell.font = title_font
                cell.fill = title_fill
                cell.alignment = title_align
                for c in range(1, col_count + 1):
                    ws.cell(row=current_row, column=c).fill = title_fill
                    _xl_apply_border(ws.cell(row=current_row, column=c), chosen["border"])
                ws.row_dimensions[current_row].height = 35
                current_row += 1

            if subtitle:
                ws.merge_cells(start_row=current_row, start_column=1, end_row=current_row, end_column=col_count)
                cell = ws.cell(row=current_row, column=1, value=subtitle)
                cell.font = sub_font
                cell.fill = sub_fill
                cell.alignment = title_align
                for c in range(1, col_count + 1):
                    ws.cell(row=current_row, column=c).fill = sub_fill
                    _xl_apply_border(ws.cell(row=current_row, column=c), chosen["border"])
                ws.row_dimensions[current_row].height = 25
                current_row += 1

            if headers:
                for col_idx, header in enumerate(headers, 1):
                    cell = ws.cell(row=current_row, column=col_idx, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = header_align
                    _xl_apply_border(cell, chosen["border"])
                ws.row_dimensions[current_row].height = 25
                current_row += 1

            for row_idx, row in enumerate(rows):
                data_row_num = row_idx + 1
                is_highlight = data_row_num in highlight_rows
                is_alt = row_idx % 2 == 1
                for col_idx, value in enumerate(row, 1):
                    cell = ws.cell(row=current_row, column=col_idx, value=value)
                    if is_highlight:
                        cell.font = highlight_font
                        cell.fill = highlight_fill
                    else:
                        cell.font = body_font
                        cell.fill = alt_fill if is_alt else white_fill
                    cell.alignment = body_align
                    _xl_apply_border(cell, chosen["border"])
                ws.row_dimensions[current_row].height = 22
                current_row += 1

            for mc in merged_cells:
                ws.merge_cells(
                    start_row=mc["start_row"], start_column=mc["start_col"],
                    end_row=mc["end_row"], end_column=mc["end_col"],
                )

            if column_widths:
                for col_idx, width in enumerate(column_widths, 1):
                    ws.column_dimensions[get_column_letter(col_idx)].width = width
            elif auto_fit:
                for col_idx in range(1, col_count + 1):
                    max_len = 0
                    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=col_idx, max_col=col_idx):
                        for cell in row:
                            if cell.value:
                                max_len = max(max_len, len(str(cell.value)))
                    ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 5, 40)

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        wb.save(output_path)
        abs_path = os.path.abspath(output_path)
        return {
            "status": "success",
            "file_path": abs_path,
            "theme": chosen["name"],
            "message": f"Excel file created with '{chosen['name']}' theme at {abs_path}",
        }
    except Exception as e:
        return {"status": "error", "message": f"Failed to create Excel file: {str(e)}"}


_THEMES = [
    {"name": "Midnight", "bg": "1E2761", "accent": "CADCFC", "title": "FFFFFF", "body": "E8E8E8", "accent2": "4A90D9"},
    {"name": "Forest", "bg": "2C5F2D", "accent": "97BC62", "title": "FFFFFF", "body": "E8F5E9", "accent2": "A5D6A7"},
    {"name": "Coral", "bg": "F96167", "accent": "FCE4EC", "title": "FFFFFF", "body": "FFF5F5", "accent2": "FFB74D"},
    {"name": "Ocean", "bg": "065A82", "accent": "1C7293", "title": "FFFFFF", "body": "E0F7FA", "accent2": "4DD0E1"},
    {"name": "Charcoal", "bg": "36454F", "accent": "F2F2F2", "title": "FFFFFF", "body": "ECEFF1", "accent2": "78909C"},
    {"name": "Berry", "bg": "6D2E46", "accent": "E1B0B6", "title": "FFFFFF", "body": "FCE4EC", "accent2": "F48FB1"},
    {"name": "Teal", "bg": "00695C", "accent": "B2DFDB", "title": "FFFFFF", "body": "E0F2F1", "accent2": "4DB6AC"},
    {"name": "Sunset", "bg": "E65100", "accent": "FFE0B2", "title": "FFFFFF", "body": "FFF3E0", "accent2": "FFB74D"},
]


def _hex_to_rgb(hex_str):
    return RGBColor.from_string(hex_str)


def _add_shape_bg(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if isinstance(text, list):
        for i, item in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.font.name = font_name
            p.space_after = Pt(font_size * 0.5)
            p.level = 0
    else:
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
    return tf


async def create_pptx(
    title: str,
    slides_json: str,
    output_path: str,
    theme: Optional[str] = None,
) -> dict:
    """
    Create a colorful, professionally styled PowerPoint presentation.

    Args:
        title: The presentation title, used on the title slide.
        slides_json: A JSON string representing a list of slides. Each slide is an object with:
            - "layout" (string): One of "title", "section", "content", "two_column", "cards", "big_number", "blank". Defaults to "content".
            - "title" (string): Slide title text.
            - "subtitle" (string): Optional subtitle.
            - "body" (string or list of strings): Body text or bullet points.
            - "body_right" (string or list of strings): Right column body for "two_column" layout.
            - "cards" (list of objects): For "cards" layout, each card has "title" and "text".
            - "number" (string): For "big_number" layout, the large number to display.
            - "number_label" (string): For "big_number" layout, label below the number.
            - "notes" (string): Optional speaker notes.
            Example: '[{"layout": "title", "title": "My Presentation", "subtitle": "A great deck"}, {"layout": "content", "title": "Key Points", "body": ["Point one", "Point two"]}]'
        output_path: File path to save the PPTX (e.g., 'presentation.pptx').
        theme: Optional theme name. Choices: Midnight, Forest, Coral, Ocean, Charcoal, Berry, Teal, Sunset. Auto-selected if not provided.

    Returns:
        Dictionary with 'status', 'file_path', 'theme', and 'message'.
    """
    try:
        slides = json.loads(slides_json)
    except json.JSONDecodeError as e:
        return {"status": "error", "message": f"Invalid JSON in slides_json: {e}"}

    try:
        import random
        chosen = None
        if theme:
            for t in _THEMES:
                if t["name"].lower() == theme.lower():
                    chosen = t
                    break
        if not chosen:
            chosen = random.choice(_THEMES)

        bg_c = _hex_to_rgb(chosen["bg"])
        accent_c = _hex_to_rgb(chosen["accent"])
        accent2_c = _hex_to_rgb(chosen["accent2"])
        title_c = _hex_to_rgb(chosen["title"])
        body_c = _hex_to_rgb(chosen["body"])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for idx, slide_data in enumerate(slides):
            layout_type = slide_data.get("layout", "content")
            slide_title = slide_data.get("title", "")
            subtitle = slide_data.get("subtitle", "")
            body = slide_data.get("body", [])
            body_right = slide_data.get("body_right", [])
            cards = slide_data.get("cards", [])
            number = slide_data.get("number", "")
            number_label = slide_data.get("number_label", "")
            notes = slide_data.get("notes", "")

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = bg_c

            if layout_type == "title":
                _add_shape_bg(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), chosen["accent2"])
                _add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.5), Inches(1.5),
                             slide_title, font_size=44, bold=True, color=title_c, alignment=PP_ALIGN.LEFT, font_name="Calibri")
                if subtitle:
                    _add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10.5), Inches(0.8),
                                 subtitle, font_size=22, color=accent_c, alignment=PP_ALIGN.LEFT)
                _add_shape_bg(slide, Inches(1.5), Inches(5.2), Inches(3), Inches(0.06), chosen["accent2"])

            elif layout_type == "section":
                _add_shape_bg(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(2.2), chosen["accent2"])
                _add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.2),
                             slide_title, font_size=36, bold=True, color=bg_c, alignment=PP_ALIGN.LEFT)
                if subtitle:
                    _add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.6),
                                 subtitle, font_size=18, color=title_c)

            elif layout_type == "content":
                _add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), chosen["accent2"])
                _add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(1.0),
                             slide_title, font_size=30, bold=True, color=bg_c, alignment=PP_ALIGN.LEFT)
                if body:
                    _add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0),
                                 body, font_size=19, color=body_c)

            elif layout_type == "two_column":
                _add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), chosen["accent2"])
                _add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(1.0),
                             slide_title, font_size=30, bold=True, color=bg_c, alignment=PP_ALIGN.LEFT)
                _add_shape_bg(slide, Inches(6.5), Inches(1.6), Inches(0.03), Inches(5.4), chosen["accent2"])
                if body:
                    _add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.3), Inches(5.0),
                                 body, font_size=17, color=body_c)
                if body_right:
                    _add_textbox(slide, Inches(7), Inches(1.8), Inches(5.3), Inches(5.0),
                                 body_right, font_size=17, color=body_c)

            elif layout_type == "cards":
                _add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(1.0),
                             slide_title, font_size=30, bold=True, color=title_c, alignment=PP_ALIGN.LEFT)
                if cards:
                    n = len(cards)
                    cols = min(n, 4)
                    card_w = (11.7 - 0.4 * (cols - 1)) / cols
                    card_h = 4.5
                    start_y = 1.8
                    for ci, card in enumerate(cards):
                        col = ci % cols
                        row_offset = (ci // cols) * (card_h + 0.4)
                        x = 0.8 + col * (card_w + 0.4)
                        y = start_y + row_offset
                        _add_shape_bg(slide, Inches(x), Inches(y), Inches(card_w), Inches(card_h), chosen["accent2"])
                        card_title = card.get("title", "")
                        card_text = card.get("text", "")
                        if card_title:
                            _add_textbox(slide, Inches(x + 0.3), Inches(y + 0.3), Inches(card_w - 0.6), Inches(0.8),
                                         card_title, font_size=20, bold=True, color=bg_c)
                        if card_text:
                            _add_textbox(slide, Inches(x + 0.3), Inches(y + 1.2), Inches(card_w - 0.6), Inches(card_h - 1.5),
                                         card_text, font_size=14, color=RGBColor(0x33, 0x33, 0x33))

            elif layout_type == "big_number":
                _add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(1.0),
                             slide_title, font_size=28, bold=True, color=title_c, alignment=PP_ALIGN.LEFT)
                if number:
                    _add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(2.5),
                                 str(number), font_size=96, bold=True, color=accent_c, alignment=PP_ALIGN.CENTER, font_name="Calibri")
                if number_label:
                    _add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(1.0),
                                 number_label, font_size=24, color=body_c, alignment=PP_ALIGN.CENTER)

            elif layout_type == "blank":
                pass

            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        prs.save(output_path)
        abs_path = os.path.abspath(output_path)
        return {
            "status": "success",
            "file_path": abs_path,
            "theme": chosen["name"],
            "message": f"PPTX created with '{chosen['name']}' theme at {abs_path}",
        }
    except Exception as e:
        return {"status": "error", "message": f"Failed to create PPTX: {str(e)}"}


async def run_cli_command(
    command: str,
    working_dir: Optional[str] = None,
    timeout: int = 60,
) -> dict:
    """
    Execute a CLI command and return the output.

    Args:
        command: The shell command to execute (e.g., 'dir', 'python script.py', 'npm install')
        working_dir: Optional working directory to run the command in. Defaults to current directory.
        timeout: Maximum execution time in seconds. Defaults to 60.

    Returns:
        Dictionary with 'stdout', 'stderr', and 'return_code' keys
    """
    try:
        process = await asyncio.create_subprocess_shell(
            command,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=working_dir,
        )

        try:
            stdout_bytes, stderr_bytes = await asyncio.wait_for(
                process.communicate(),
                timeout=timeout,
            )
        except asyncio.TimeoutError:
            process.kill()
            await process.communicate()
            return {
                "stdout": "",
                "stderr": f"Command timed out after {timeout} seconds",
                "return_code": "-1",
            }

        stdout = stdout_bytes.decode("utf-8", errors="replace")
        stderr = stderr_bytes.decode("utf-8", errors="replace")

        return {
            "stdout": stdout[:10000],
            "stderr": stderr[:10000],
            "return_code": str(process.returncode),
        }

    except Exception as e:
        return {
            "stdout": "",
            "stderr": f"Failed to execute command: {str(e)}",
            "return_code": "-1",
        }


create_pdf_tool = FunctionTool(create_pdf)
create_excel_tool = FunctionTool(create_excel)
create_pptx_tool = FunctionTool(create_pptx)
run_cli_command_tool = FunctionTool(run_cli_command)
